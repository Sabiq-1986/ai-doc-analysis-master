"""
PowerPoint Parser - python-pptx with comprehensive extraction.
Slide text with font-size heading detection, speaker notes, tables with merged cells,
charts with data extraction, grouped shapes (recursive), embedded images with OCR, comments.
"""
import io
import logging
from pathlib import Path
from typing import List, Optional

from langchain_core.documents import Document

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXParser:
    """Comprehensive PowerPoint parser with full feature extraction."""

    def __init__(self, image_parser=None):
        self.image_parser = image_parser

    def parse(self, file_path: str) -> List[Document]:
        """Parse PowerPoint file. Returns one Document per slide."""
        filename = Path(file_path).name

        if not PPTX_AVAILABLE:
            logger.error("[PPTX] python-pptx not installed")
            return []

        documents = []

        try:
            prs = Presentation(file_path)

            # Presentation properties
            props = self._extract_properties(prs)
            if props:
                documents.append(Document(
                    page_content=props,
                    metadata={"source": filename, "type": "pptx", "section": "properties"}
                ))

            total_slides = len(prs.slides)

            for slide_idx, slide in enumerate(prs.slides):
                slide_num = slide_idx + 1
                slide_parts = []

                # Slide title
                title = self._get_slide_title(slide)
                if title:
                    slide_parts.append(f"# {title}")

                # Slide layout
                try:
                    layout_name = slide.slide_layout.name
                    slide_parts.append(f"[Layout: {layout_name}]")
                except Exception:
                    pass

                # Text from all shapes (with font-size heading detection)
                text = self._extract_slide_text(slide)
                if text:
                    slide_parts.append(text)

                # Tables
                tables = self._extract_tables(slide)
                if tables:
                    slide_parts.append(tables)

                # Charts
                charts = self._extract_charts(slide)
                if charts:
                    slide_parts.append(charts)

                # Speaker notes
                notes = self._extract_speaker_notes(slide)
                if notes:
                    slide_parts.append(f"\n[Speaker Notes]\n{notes}")

                # Comments
                comments = self._extract_comments(slide, slide_num)
                if comments:
                    slide_parts.append(comments)

                # Embedded images
                image_docs = self._extract_images(slide, slide_num, filename)
                documents.extend(image_docs)

                # Create slide document
                slide_content = '\n\n'.join(slide_parts)
                if slide_content.strip():
                    documents.append(Document(
                        page_content=slide_content.strip(),
                        metadata={
                            "source": filename,
                            "type": "pptx",
                            "slide": slide_num,
                            "total_slides": total_slides,
                            "title": title or f"Slide {slide_num}",
                        }
                    ))

        except Exception as e:
            logger.error(f"[PPTX] Error parsing {filename}: {e}")

        return documents

    def _extract_properties(self, prs) -> Optional[str]:
        """Extract presentation properties."""
        try:
            props = prs.core_properties
            lines = ["[Presentation Properties]"]

            if props.title:
                lines.append(f"Title: {props.title}")
            if props.author:
                lines.append(f"Author: {props.author}")
            if props.subject:
                lines.append(f"Subject: {props.subject}")
            if props.keywords:
                lines.append(f"Keywords: {props.keywords}")
            if props.created:
                lines.append(f"Created: {props.created}")
            if props.modified:
                lines.append(f"Modified: {props.modified}")

            lines.append(f"Slides: {len(prs.slides)}")

            # Slide dimensions
            try:
                w = prs.slide_width
                h = prs.slide_height
                if w and h:
                    lines.append(f"Dimensions: {w/914400:.1f}\" x {h/914400:.1f}\"")
            except Exception:
                pass

            if len(lines) > 2:
                return '\n'.join(lines)
        except Exception:
            pass
        return None

    def _get_slide_title(self, slide) -> Optional[str]:
        """Get slide title from title placeholder."""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
        except Exception:
            pass

        # Fallback: look for the largest text
        try:
            largest_size = 0
            largest_text = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size > largest_size:
                                if run.text.strip():
                                    largest_size = run.font.size
                                    largest_text = run.text.strip()
            return largest_text
        except Exception:
            pass

        return None

    def _extract_slide_text(self, slide) -> str:
        """Extract text from all shapes with font-size-based heading detection."""
        text_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = self._extract_shape_text(shape)
                if shape_text:
                    text_parts.append(shape_text)

            # Handle grouped shapes recursively
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = self._extract_group_text(shape)
                if group_text:
                    text_parts.append(group_text)

        return '\n'.join(text_parts)

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape with font-size awareness."""
        if not shape.has_text_frame:
            return ""

        lines = []
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect heading by font size
            max_size = 0
            has_bold = False
            for run in para.runs:
                if run.font.size:
                    size_pt = run.font.size / 12700  # EMU to points
                    max_size = max(max_size, size_pt)
                if run.font.bold:
                    has_bold = True

            if max_size >= 24:
                lines.append(f"## {text}")
            elif max_size >= 18 or has_bold:
                lines.append(f"**{text}**")
            else:
                # Check for bullet points
                level = para.level if hasattr(para, 'level') else 0
                if level > 0:
                    indent = "  " * level
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)

        return '\n'.join(lines)

    def _extract_group_text(self, group_shape) -> str:
        """Recursively extract text from grouped shapes."""
        text_parts = []

        try:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    shape_text = self._extract_shape_text(shape)
                    if shape_text:
                        text_parts.append(shape_text)

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = self._extract_group_text(shape)
                    if group_text:
                        text_parts.append(group_text)
        except Exception:
            pass

        return '\n'.join(text_parts)

    def _extract_tables(self, slide) -> str:
        """Extract tables with merged cell handling."""
        table_texts = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = []

                for row_idx, row in enumerate(table.rows):
                    cells = []
                    for col_idx, cell in enumerate(row.cells):
                        text = cell.text.strip()

                        # Check for merged cells
                        if cell.is_merge_origin:
                            span_h = cell.span_width if hasattr(cell, 'span_width') else 1
                            span_v = cell.span_height if hasattr(cell, 'span_height') else 1
                            if span_h > 1 or span_v > 1:
                                text = f"{text} [merged]"

                        cells.append(text)

                    if any(cells):
                        rows.append(' | '.join(cells))

                if rows:
                    table_texts.append("[Table]\n" + '\n'.join(rows))

        return '\n\n'.join(table_texts)

    def _extract_charts(self, slide) -> str:
        """Extract chart data including type, title, series, and categories."""
        chart_texts = []

        for shape in slide.shapes:
            if shape.has_chart:
                try:
                    chart = shape.chart
                    parts = []

                    # Chart type
                    chart_type = str(chart.chart_type) if chart.chart_type else "Unknown"
                    parts.append(f"[Chart: {chart_type}]")

                    # Chart title
                    if chart.has_title and chart.chart_title:
                        try:
                            title_text = chart.chart_title.text_frame.text
                            if title_text:
                                parts.append(f"Title: {title_text}")
                        except Exception:
                            pass

                    # Series data
                    try:
                        for s_idx, series in enumerate(chart.series):
                            series_name = series.name if hasattr(series, 'name') else f"Series {s_idx+1}"
                            values = []
                            try:
                                for point in series.values:
                                    if point is not None:
                                        values.append(str(point))
                            except Exception:
                                pass

                            if values:
                                parts.append(f"  {series_name}: {', '.join(values[:20])}")
                    except Exception:
                        pass

                    # Categories
                    try:
                        categories = []
                        plot = chart.plots[0] if chart.plots else None
                        if plot:
                            for cat in plot.categories:
                                if cat:
                                    categories.append(str(cat))
                        if categories:
                            parts.append(f"  Categories: {', '.join(categories[:20])}")
                    except Exception:
                        pass

                    if len(parts) > 1:
                        chart_texts.append('\n'.join(parts))

                except Exception as e:
                    logger.debug(f"[PPTX] Chart extraction failed: {e}")

        return '\n\n'.join(chart_texts)

    def _extract_speaker_notes(self, slide) -> Optional[str]:
        """Extract speaker notes from slide."""
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                text = notes_slide.notes_text_frame.text.strip()
                if text:
                    return text
        except Exception:
            pass
        return None

    def _extract_comments(self, slide, slide_num: int) -> Optional[str]:
        """Extract comments from slide."""
        try:
            comments = []
            # Access comments through XML
            slide_element = slide._element
            for comment in slide_element.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cm'):
                author = comment.get('authorId', '')
                text_elem = comment.find('{http://schemas.openxmlformats.org/presentationml/2006/main}text')
                if text_elem is not None and text_elem.text:
                    comments.append(f"Comment: {text_elem.text.strip()}")

            if comments:
                return '\n'.join(comments)
        except Exception:
            pass
        return None

    def _extract_images(self, slide, slide_num: int, filename: str) -> List[Document]:
        """Extract and OCR embedded images from slide."""
        if not self.image_parser:
            return []

        documents = []
        img_idx = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_data = image.blob
                    img_idx += 1

                    if len(img_data) < 1000:
                        continue

                    result = self.image_parser.extract_text(
                        img_data,
                        filename=f"{filename}_s{slide_num}_img{img_idx}",
                        mode="auto"
                    )

                    text = result.get("text", "")
                    if text and len(text.strip()) > 10:
                        documents.append(Document(
                            page_content=text,
                            metadata={
                                "source": filename,
                                "type": "pptx_image",
                                "slide": slide_num,
                                "image_index": img_idx,
                            }
                        ))

                except Exception as e:
                    logger.debug(f"[PPTX] Image extraction failed slide {slide_num}: {e}")

        return documents
